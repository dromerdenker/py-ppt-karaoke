import os
import random
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

NETT_SLIDES_PER_DECK = 5

def add_text_slide(prs, slammer, supporter=""):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slammer
    subtitle.text = supporter


def add_random_slide(prs):
    img_path = random_slide_file_name()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(0)
    top = Inches(0)
    slide = slide.shapes.add_picture(
        img_path, Inches(0.1), Inches(0.1),
        width=Inches(9), height=Inches(5))


def random_slide_file_name():
    p = Path(__file__).parent.parent / 'assets'
    files = os.listdir(p)
    return str(p / random.choice(files))


def generate_one_deck(deckid):
    ppt_path = Path(__file__).parent.parent / 'generated' / f'deck_{deckid:02}.pptx'

    prs = Presentation()
    add_text_slide(prs, "Welcome to ppt-karaoke", f"Deck-{deckid:02}")

    for sl in range(NETT_SLIDES_PER_DECK):
        add_random_slide(prs)
    add_text_slide(prs, "Thank you!", "Who is next? :-)")

    prs.save(str(ppt_path))

if __name__ == '__main__':
    for deck in range(5):
        generate_one_deck(deck)
